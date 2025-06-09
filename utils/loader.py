from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import requests
from langchain.docstore.document import Document as LangDoc

def extract_text_from_pdf(file):
    return "\n".join([p.extract_text() for p in PdfReader(file).pages])

def extract_text_from_docx(file):
    return "\n".join([para.text for para in Document(file).paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def extract_text_from_url(url):
    try:
        html = requests.get(url).text
        return BeautifulSoup(html, "html.parser").get_text()
    except:
        return ""

def load_documents(uploaded_files, url_input):
    docs = []
    for file in uploaded_files:
        if file.name.endswith(".pdf"):
            text = extract_text_from_pdf(file)
        elif file.name.endswith(".docx"):
            text = extract_text_from_docx(file)
        elif file.name.endswith(".pptx"):
            text = extract_text_from_pptx(file)
        else:
            continue
        docs.append(LangDoc(page_content=text, metadata={"source": file.name}))
    if url_input:
        docs.append(LangDoc(page_content=extract_text_from_url(url_input), metadata={"source": url_input}))
    return docs
